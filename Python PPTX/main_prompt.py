import collections 
import collections.abc
from pptx import Presentation

def make(title, name, arr):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # first slide - title

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = name

    # second slide - content

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title

    tf = body_shape.text_frame
    tf.text = 'Bullet points added:'

    for txt in arr:
        p = tf.add_paragraph()
        p.text = txt
        p.level = 1

    prs.save('test.pptx')

if __name__ == "__main__":
    print("Creating a PowerPoint...")
    title = input("Title of presentation: ")
    name = input("Name: ")
    name = "By " + name
    numBullets = int(input("How many bullet points?: "))
    arrBullets = [None] * numBullets
    for i in range(numBullets):
        tmp = input("Bullet point " + str(i + 1) + ": ")
        arrBullets[i] = tmp
    make(title, name, arrBullets)
    print("PowerPoint created as test.pptx!")